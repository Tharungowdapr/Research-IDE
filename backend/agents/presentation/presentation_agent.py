"""
Presentation Agent — Generates slide content and PPTX presentation decks.
"""

import json
import re
import io
from typing import Dict, List, Any, Optional
from core.llm_client import LLMClient

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


SLIDE_SYSTEM = "You are an expert at creating research presentation content. Return ONLY valid JSON."

SLIDE_PROMPT = """Generate presentation slide content for this research project.

Title: {title}
Domain: {domain}
Idea: {description}
Approach: {approach}
Novelty: {novelty}
Overview: {overview}
Metrics: {metrics}
Datasets: {datasets}

Related Papers:
{related_papers}

Return EXACTLY this JSON structure:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "subtitle": "Optional subtitle",
      "bullets": ["bullet point 1", "bullet point 2", "bullet point 3"],
      "notes": "Speaker notes for this slide"
    }}
  ]
}}

Generate 8-12 slides covering: title, problem statement, motivation, related work, methodology, experimental setup, expected results, timeline, conclusion, references."""


async def generate_slide_content(
    idea: Dict,
    papers: List[Dict],
    plan: Dict,
    intent: Dict,
    llm: LLMClient,
) -> Dict:
    """Generate structured slide content using the LLM."""
    domain = ", ".join(intent.get("domain", ["AI/ML"]))
    related = _format_papers_short(papers[:5])
    metrics = plan.get("evaluation_metrics", [])
    datasets = [d.get("name", "") for d in plan.get("datasets", [])]

    prompt = SLIDE_PROMPT.format(
        title=idea.get("title", "Research Project"),
        domain=domain,
        description=idea.get("description", ""),
        approach=idea.get("approach", ""),
        novelty=idea.get("novelty", ""),
        overview=plan.get("overview", ""),
        metrics=", ".join(str(m) for m in metrics if m) or "accuracy, F1",
        datasets=", ".join(str(d) for d in datasets if d) or "standard benchmarks",
        related_papers=related,
    )

    try:
        raw = await llm.complete(prompt, system=SLIDE_SYSTEM, json_mode=True)
        result = _parse_json(raw)
        if "slides" in result and len(result["slides"]) > 0:
            return result
        return _fallback_slides(idea, plan)
    except Exception as e:
        print(f"Slide content generation error: {e}")
        return _fallback_slides(idea, plan)


def export_pptx(slides_data: Dict, output_path: Optional[str] = None) -> bytes:
    """Generate a PPTX file from slide content."""
    if not PPTX_AVAILABLE:
        return json.dumps(slides_data).encode()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
    ACCENT = RGBColor(0x81, 0x81, 0xFF)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

    slides = slides_data.get("slides", [])

    for i, slide_data in enumerate(slides):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        bullets = slide_data.get("bullets", [])

        left = Inches(0.8)
        top = Inches(0.6)
        width = Inches(11.7)

        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        if subtitle:
            sub_box = slide.shapes.add_textbox(left, Inches(1.5), width, Inches(0.5))
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(18)
            sp.font.color.rgb = LIGHT_GRAY

        if bullets:
            bullet_top = Inches(2.2) if subtitle else Inches(1.7)
            bullet_box = slide.shapes.add_textbox(left, bullet_top, width, Inches(4.5))
            btf = bullet_box.text_frame
            btf.word_wrap = True

            for j, bullet in enumerate(bullets):
                if j == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()

                bp.text = f"• {bullet}"
                bp.font.size = Pt(16)
                bp.font.color.rgb = WHITE
                bp.space_after = Pt(8)

        if i == 0:
            footer_box = slide.shapes.add_textbox(left, Inches(6.8), width, Inches(0.4))
            ftf = footer_box.text_frame
            fp = ftf.paragraphs[0]
            fp.text = "ResearchIDE — AI-Powered Research Assistant"
            fp.font.size = Pt(10)
            fp.font.color.rgb = LIGHT_GRAY
            fp.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _fallback_slides(idea: Dict, plan: Dict) -> Dict:
    title = idea.get("title", "Research Project")
    return {
        "slides": [
            {
                "title": title,
                "subtitle": "A ResearchIDE Project",
                "bullets": ["Venue: Conference / Journal Submission", "Domain: AI/ML Research"],
                "notes": "Introduce the research project and its significance.",
            },
            {
                "title": "Problem Statement",
                "subtitle": "",
                "bullets": [
                    idea.get("description", "Define the research problem here."),
                    f"Novelty: {idea.get('novelty', 'Novel approach')}",
                ],
                "notes": "Motivate the problem and explain why it matters.",
            },
            {
                "title": "Related Work",
                "subtitle": "",
                "bullets": ["Reviewed key papers in the domain", "Identified gaps in existing approaches"],
                "notes": "Summarize the literature review.",
            },
            {
                "title": "Proposed Approach",
                "subtitle": "",
                "bullets": [
                    f"Methodology: {idea.get('approach', 'Novel methodology')}",
                    "Designed to address identified gaps",
                ],
                "notes": "Explain the technical approach.",
            },
            {
                "title": "Experimental Setup",
                "subtitle": "",
                "bullets": [
                    f"Datasets: {', '.join(d.get('name', '') for d in plan.get('datasets', [])) or 'Standard benchmarks'}",
                    f"Metrics: {', '.join(plan.get('evaluation_metrics', [])) or 'Accuracy, F1'}",
                ],
                "notes": "Describe evaluation methodology.",
            },
            {
                "title": "Timeline",
                "subtitle": "",
                "bullets": [
                    f"{p.get('name', '')}: {p.get('duration', 'TBD')}"
                    for p in plan.get("phases", [])
                ],
                "notes": "Walk through the project timeline.",
            },
            {
                "title": "Expected Contributions",
                "subtitle": "",
                "bullets": [
                    "Novel methodology advancing the state of the art",
                    "Comprehensive experimental evaluation",
                    "Open-source code and reproducible results",
                ],
                "notes": "Highlight the expected impact.",
            },
            {
                "title": "Conclusion & Next Steps",
                "subtitle": "",
                "bullets": [
                    "Clear problem definition and approach",
                    "Systematic evaluation plan",
                    "Next: data collection and implementation",
                ],
                "notes": "Summarize and outline next steps.",
            },
        ],
        "_fallback": True,
    }


def _format_papers_short(papers: List[Dict]) -> str:
    return "\n".join(
        f"- {p.get('title', '')} ({p.get('year', '')})"
        for p in papers if p.get("title")
    )


def _parse_json(raw: str) -> Dict:
    clean = re.sub(r"```(?:json)?\s*", "", raw).strip().rstrip("`").strip()
    s, e = clean.find("{"), clean.rfind("}") + 1
    return json.loads(clean[s:e]) if s != -1 and e > s else {}
