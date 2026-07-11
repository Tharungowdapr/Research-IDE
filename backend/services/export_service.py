"""
Document Export Utilities — DOCX, PDF (fpdf2), and PPTX generation for IEEE-format papers.
"""

import io
import re
from typing import Dict, List
from docx import Document
from docx.shared import Pt, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml


def generate_docx(report: Dict) -> bytes:
    """Generate IEEE-format DOCX from report dict."""
    doc = Document()

    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(1.9)
    section.bottom_margin = Cm(2.54)
    section.left_margin = Cm(1.9)
    section.right_margin = Cm(1.9)

    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_para.space_before = Pt(0)
    title_para.space_after = Pt(12)
    title_run = title_para.add_run(report.get("title", "Research Paper"))
    title_run.bold = True
    title_run.font.size = Pt(24)
    title_run.font.name = "Times New Roman"

    authors = report.get("authors", ["Author"])
    authors_para = doc.add_paragraph()
    authors_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    authors_para.space_before = Pt(0)
    authors_para.space_after = Pt(12)
    auth_run = authors_para.add_run(", ".join(authors))
    auth_run.italic = True
    auth_run.font.size = Pt(11)
    auth_run.font.name = "Times New Roman"

    affiliations = report.get("affiliations", [])
    if affiliations:
        aff_para = doc.add_paragraph()
        aff_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        aff_para.space_before = Pt(0)
        aff_para.space_after = Pt(4)
        aff_run = aff_para.add_run("; ".join(affiliations))
        aff_run.font.size = Pt(9)
        aff_run.font.name = "Times New Roman"

    emails = report.get("emails", [])
    if emails:
        em_para = doc.add_paragraph()
        em_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        em_para.space_before = Pt(0)
        em_para.space_after = Pt(8)
        em_run = em_para.add_run(", ".join(emails))
        em_run.italic = True
        em_run.font.size = Pt(8)
        em_run.font.name = "Times New Roman"

    abs_label = doc.add_paragraph()
    abs_label.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    abs_label.space_before = Pt(0)
    abs_label.space_after = Pt(6)
    label_run = abs_label.add_run("Abstract—")
    label_run.italic = True
    label_run.font.size = Pt(9)
    label_run.font.name = "Times New Roman"
    abs_run = abs_label.add_run(report.get("abstract", ""))
    abs_run.font.size = Pt(9)
    abs_run.font.name = "Times New Roman"

    kw_para = doc.add_paragraph()
    kw_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    kw_para.space_before = Pt(0)
    kw_para.space_after = Pt(12)
    kw_label = kw_para.add_run("Index Terms—")
    kw_label.italic = True
    kw_label.font.size = Pt(9)
    kw_label.font.name = "Times New Roman"
    kw_text = kw_para.add_run(", ".join(report.get("keywords", [])))
    kw_text.font.size = Pt(9)
    kw_text.font.name = "Times New Roman"

    new_section = doc.add_section(WD_SECTION.CONTINUOUS)
    new_section.page_width = Cm(21)
    new_section.page_height = Cm(29.7)
    new_section.top_margin = Cm(1.9)
    new_section.bottom_margin = Cm(2.54)
    new_section.left_margin = Cm(1.9)
    new_section.right_margin = Cm(1.9)
    sect_pr = new_section._sectPr
    cols_xml = f'<w:cols {nsdecls("w")} w:num="2" w:space="720"/>'
    sect_pr.append(parse_xml(cols_xml))

    for sec in report.get("sections", []):
        heading = sec.get("heading", "")
        content = sec.get("content", "")

        h_para = doc.add_paragraph()
        h_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        h_para.space_before = Pt(12)
        h_para.space_after = Pt(6)
        h_run = h_para.add_run(heading.upper())
        h_run.bold = True
        h_run.font.size = Pt(10)
        h_run.font.name = "Times New Roman"

        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        for para_text in paragraphs:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            p.space_before = Pt(0)
            p.space_after = Pt(6)
            run = p.add_run(para_text)
            run.font.size = Pt(10)
            run.font.name = "Times New Roman"

    ack = report.get("acknowledgements", "")
    if ack:
        ack_h = doc.add_paragraph()
        ack_h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ack_h.space_before = Pt(12)
        ack_h.space_after = Pt(6)
        ack_hr = ack_h.add_run("ACKNOWLEDGEMENTS")
        ack_hr.bold = True
        ack_hr.font.size = Pt(10)
        ack_hr.font.name = "Times New Roman"
        ack_p = doc.add_paragraph()
        ack_p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        ack_p.space_after = Pt(6)
        ack_r = ack_p.add_run(ack)
        ack_r.font.size = Pt(9)
        ack_r.font.name = "Times New Roman"

    refs = report.get("references", [])
    if refs:
        ref_h = doc.add_paragraph()
        ref_h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ref_h.space_before = Pt(12)
        ref_h.space_after = Pt(6)
        ref_hr = ref_h.add_run("REFERENCES")
        ref_hr.bold = True
        ref_hr.font.size = Pt(10)
        ref_hr.font.name = "Times New Roman"

        for ref in refs:
            ref_p = doc.add_paragraph()
            ref_p.space_before = Pt(0)
            ref_p.space_after = Pt(3)
            if isinstance(ref, dict):
                ref_id = ref.get("id", "")
                authors_str = ref.get("authors", "Unknown")
                title_str = ref.get("title", "")
                venue = ref.get("venue", "")
                year = ref.get("year", "")
                ref_text = f"[{ref_id}] {authors_str}, \"{title_str},\" {venue}, {year}."
            else:
                ref_text = str(ref)
            run = ref_p.add_run(ref_text)
            run.font.size = Pt(8)
            run.font.name = "Times New Roman"

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _sanitize_pdf(text: str) -> str:
    """Replace Unicode chars that fpdf2 can't render."""
    if not text:
        return ""
    replacements = {
        "\u2014": " -- ", "\u2013": " - ", "\u2018": "'", "\u2019": "'",
        "\u201c": '"', "\u201d": '"', "\u2026": "...", "\u2022": "*",
        "\u00a0": " ", "\u2192": "->", "\u2190": "<-", "\u2264": "<=",
        "\u2265": ">=", "\u00d7": "x", "\u00f7": "/",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def generate_pdf(report: Dict) -> bytes:
    """Generate a proper IEEE-format two-column PDF with Times New Roman."""
    from fpdf import FPDF
    import os

    # ── IEEE layout constants (A4 in mm) ──
    PAGE_W = 210
    PAGE_H = 297
    MARGIN_L = 15
    MARGIN_R = 15
    MARGIN_TOP = 15
    MARGIN_BOT = 15
    COL_GAP = 8
    COL_W = (PAGE_W - MARGIN_L - MARGIN_R - COL_GAP) / 2  # ~86mm each
    BODY_TOP = MARGIN_TOP  # updated after header

    pdf = FPDF(orientation='P', unit='mm', format='A4')
    pdf.set_auto_page_break(auto=False)
    pdf.set_margins(MARGIN_L, MARGIN_TOP, MARGIN_R)

    # Register Times New Roman TTF
    font_dir = r"C:\Windows\Fonts"
    fonts = {
        "":   os.path.join(font_dir, "times.ttf"),
        "B":  os.path.join(font_dir, "timesbd.ttf"),
        "I":  os.path.join(font_dir, "timesi.ttf"),
        "BI": os.path.join(font_dir, "timesbi.ttf"),
    }
    for style, path in fonts.items():
        if os.path.exists(path):
            pdf.add_font("TNR", style, path)

    def _col_x(col: int) -> float:
        return MARGIN_L + col * (COL_W + COL_GAP)

    def _bottom() -> float:
        return PAGE_H - MARGIN_BOT

    def _write_header():
        """Single-column header: title, authors+affiliations, abstract, keywords."""
        pdf.add_page()

        # Title
        pdf.set_font("TNR", "B", 24)
        pdf.set_x(MARGIN_L)
        pdf.multi_cell(PAGE_W - MARGIN_L - MARGIN_R, 10, _sanitize_pdf(
            report.get("title", "Research Paper")), align="C")
        pdf.ln(6)

        # Authors with affiliation markers
        authors = report.get("authors", ["Author"])
        affiliations = report.get("affiliations", [])
        emails = report.get("emails", [])

        # Author names (centered, italic)
        sup_digits = ["\u00b9", "\u00b2", "\u00b3", "\u2074", "\u2075", "\u2076", "\u2077", "\u2078", "\u2079"]
        author_line = ""
        for i, name in enumerate(authors):
            if i > 0:
                author_line += ",  "
            sup = sup_digits[i] if i < len(sup_digits) and len(authors) > 1 else ""
            author_line += f"{_sanitize_pdf(name)}{sup}"
        pdf.set_font("TNR", "I", 11)
        pdf.set_x(MARGIN_L)
        pdf.multi_cell(PAGE_W - MARGIN_L - MARGIN_R, 6, author_line, align="C")
        pdf.ln(2)

        # Affiliations
        if affiliations:
            sup_digits = ["\u00b9", "\u00b2", "\u00b3", "\u2074", "\u2075", "\u2076", "\u2077", "\u2078", "\u2079"]
            for i, aff in enumerate(affiliations):
                marker = sup_digits[i] if i < len(sup_digits) and len(affiliations) > 1 else ""
                pdf.set_font("TNR", "", 9)
                pdf.set_x(MARGIN_L)
                pdf.multi_cell(PAGE_W - MARGIN_L - MARGIN_R, 4.5,
                               f"{marker}{_sanitize_pdf(aff)}", align="C")
            pdf.ln(1)

        # Emails
        if emails:
            email_str = ", ".join(_sanitize_pdf(e) for e in emails)
            pdf.set_font("TNR", "I", 8)
            pdf.set_x(MARGIN_L)
            pdf.cell(PAGE_W - MARGIN_L - MARGIN_R, 4, f"Corresponding author: {email_str}",
                     align="C", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(3)

        # Abstract
        pdf.set_font("TNR", "B", 9)
        pdf.set_x(MARGIN_L)
        abs_label_w = pdf.get_string_width("Abstract - ") + 2
        pdf.cell(abs_label_w, 5, "Abstract - ", new_x="END", new_y="TOP")
        pdf.set_font("TNR", "", 9)
        remaining_w = PAGE_W - MARGIN_L - MARGIN_R - abs_label_w
        pdf.multi_cell(remaining_w, 4.5, _sanitize_pdf(
            report.get("abstract", "")), align="J")
        pdf.ln(2)

        # Index Terms
        keywords = report.get("keywords", [])
        if keywords:
            pdf.set_font("TNR", "B", 9)
            pdf.set_x(MARGIN_L)
            kw_label_w = pdf.get_string_width("Index Terms - ") + 2
            pdf.cell(kw_label_w, 5, "Index Terms - ", new_x="END", new_y="TOP")
            pdf.set_font("TNR", "I", 9)
            remaining_w = PAGE_W - MARGIN_L - MARGIN_R - kw_label_w
            pdf.multi_cell(remaining_w, 5, _sanitize_pdf(", ".join(keywords)), align="L")
            pdf.ln(2)

        # Horizontal rule
        y_rule = pdf.get_y() + 1
        pdf.line(MARGIN_L, y_rule, PAGE_W - MARGIN_R, y_rule)
        pdf.set_y(y_rule + 3)

    def _add_page_with_cols():
        pdf.add_page()
        pdf.set_y(MARGIN_TOP)
        pdf.set_x(_col_x(0))

    # ── Track column positions ──
    col_y = [0.0, 0.0]
    active_col = 0

    def _switch_col():
        nonlocal active_col
        if active_col == 0:
            active_col = 1
            pdf.set_y(col_y[1])
            pdf.set_x(_col_x(1))
        else:
            _add_page_with_cols()
            active_col = 0
            col_y[0] = MARGIN_TOP
            col_y[1] = MARGIN_TOP
            pdf.set_y(MARGIN_TOP)
            pdf.set_x(_col_x(0))

    def _check_space(needed_mm: float) -> bool:
        if pdf.get_y() + needed_mm > _bottom():
            _switch_col()
            return False
        return True

    def _word_wrap(text: str, font_style: str, font_size: float) -> list:
        """Manually word-wrap text into lines that fit within COL_W."""
        pdf.set_font("TNR", font_style, font_size)
        words = _sanitize_pdf(text).split()
        lines = []
        current = ""
        for word in words:
            test = f"{current} {word}".strip() if current else word
            if pdf.get_string_width(test) <= COL_W:
                current = test
            else:
                if current:
                    lines.append(current)
                current = word
        if current:
            lines.append(current)
        return lines if lines else [""]

    def _col_write_line(line: str, font_style: str, font_size: float,
                        align: str = "L", line_h: float = 4.5):
        """Write a single line in the current column."""
        if pdf.get_y() + line_h > _bottom():
            _switch_col()
        pdf.set_font("TNR", font_style, font_size)
        pdf.set_x(_col_x(active_col))
        safe = "L" if align == "J" else align
        pdf.cell(COL_W, line_h, line, align=safe, new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(_col_x(active_col))

    def _col_multi_cell(text: str, font_style: str = "", font_size: float = 10,
                        align: str = "J", line_h: float = 4.5):
        """Write paragraph text within current column, splitting across columns/pages as needed."""
        lines = _word_wrap(text, font_style, font_size)
        for line in lines:
            _col_write_line(line, font_style, font_size, "L", line_h)
        col_y[active_col] = pdf.get_y()

    def _col_heading(numeral: str, title: str):
        """Write an IEEE section heading within the column."""
        full = f"{numeral} {_sanitize_pdf(title.upper())}"
        _check_space(12)
        _col_write_line(full, "B", 10, "L", 6)
        pdf.ln(1)
        col_y[active_col] = pdf.get_y()

    # ── Generate PDF ──
    _write_header()

    BODY_TOP = pdf.get_y()
    col_y = [BODY_TOP, BODY_TOP]
    active_col = 0
    pdf.set_x(_col_x(0))

    sections = report.get("sections", [])
    numerals = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X",
                "XI", "XII", "XIII", "XIV", "XV", "XVI", "XVII", "XVIII", "XIX", "XX"]

    for i, sec in enumerate(sections):
        heading = sec.get("heading", "")
        content = sec.get("content", "")

        numeral = numerals[i] if i < len(numerals) else str(i + 1)

        if heading.strip():
            # Strip existing numeral prefix like "I. " or "1. " if present
            clean_heading = re.sub(r'^[IVXLC]+[\.\s]+', '', heading.strip())
            clean_heading = re.sub(r'^\d+[\.\s]+', '', clean_heading)
            if clean_heading.strip():
                _col_heading(numeral, clean_heading)
            else:
                _col_heading(numeral, heading)

        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        for para_text in paragraphs:
            _col_multi_cell(para_text)
            pdf.ln(2.5)
            col_y[active_col] = pdf.get_y()

    # ── Acknowledgements ──
    ack = report.get("acknowledgements", "")
    if ack:
        _check_space(20)
        _col_write_line("ACKNOWLEDGMENT", "B", 10, "L", 6)
        pdf.ln(1)
        _col_multi_cell(ack)
        pdf.ln(2)
        col_y[active_col] = pdf.get_y()

    # ── References ──
    refs = report.get("references", [])
    if refs:
        _check_space(20)
        _col_write_line("REFERENCES", "B", 10, "L", 6)
        pdf.ln(1)
        col_y[active_col] = pdf.get_y()

        for ref in refs:
            if isinstance(ref, dict):
                ref_id = ref.get("id", "")
                authors_str = ref.get("authors", "Unknown")
                title_str = ref.get("title", "")
                venue = ref.get("venue", "")
                year = ref.get("year", "")
                ref_text = f"[{ref_id}] {authors_str}, \"{title_str},\" {venue}, {year}."
            else:
                ref_text = str(ref)
            _col_multi_cell(ref_text, font_size=8, line_h=3.8)
            pdf.ln(0.8)
            col_y[active_col] = pdf.get_y()

    # ── Page numbers (bottom center, all pages) ──
    for pg in range(1, pdf.pages_count + 1):
        pdf.page = pg
        pdf.set_y(-10)
        pdf.set_font("TNR", "", 9)
        pdf.cell(0, 5, str(pg), align="C")

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf.getvalue()


def generate_pptx(presentation: dict) -> bytes:
    """Generate PPTX from presentation slide data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptxPt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = presentation.get("slides", [])

    for i, slide_data in enumerate(slides_data):
        layout = prs.slide_layouts[6] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0F)

        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        bullets = slide_data.get("bullets", [])

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xE8, 0xE8, 0xF0)

        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = PptxPt(18)
            p2.font.color.rgb = RGBColor(0x88, 0x88, 0xA8)

        if bullets:
            txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            for j, bullet in enumerate(bullets):
                p3 = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
                p3.text = f"\u2022  {bullet}"
                p3.font.size = PptxPt(16)
                p3.font.color.rgb = RGBColor(0x88, 0x88, 0xA8)
                p3.space_after = PptxPt(8)

        notes_text = slide_data.get("notes", "")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _esc(text: str) -> str:
    """Escape HTML entities."""
    return (text.replace("&", "&amp;").replace("<", "&lt;")
            .replace(">", "&gt;").replace('"', "&quot;"))
